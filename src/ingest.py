""" Document ingestion module. """

from dataclasses import dataclass, replace
from pathlib import Path

from openai import OpenAI

import os
import time
import argparse

from typing import Iterator, Sequence

from pptx import Presentation
from pypdf import PdfReader

import math
import logging

from src.config import Settings, get_settings
from src.vector_store import DocumentChunk, get_vector_store
from src.openai_client import OpenAIClient, OpenAIClientConfig

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format="%(levelname)s %(name)s: %(message)s")

# Variables
SUPPORTED_FILE_TYPES = {".txt", ".md", ".pdf", ".pptx"}
DEFAULT_EMBED_BATCH_SIZE = get_settings().embed_batch_size or 64

# Represents a document that has been loaded from a file
# with its file path and content. 
# It's immutable to ensure data integrity.
@dataclass(frozen=True) 
class LoadedDocument:
    path: Path
    content: str


# Loads all text documents from the specified root directory.
def load_documents(root: Path) -> list[LoadedDocument]:
    documents: list[LoadedDocument] = []

    # Check if the root path exists 
    if not root.exists():
        logger.warning(f"Documents Directory {root} does not exist.")
        # print(f"Documents Directory {root} does not exist.")
        return documents
    
    # Iterate and check if each file is supported or not
    for file_path in root.rglob('*'):
        if file_path.is_file() and file_path.suffix.lower() not in SUPPORTED_FILE_TYPES:
            logger.info(f"Skipping unsupported file type: {file_path}")
            continue
        # print(f"Processing file: {file_path}")
        logger.info(f"Processing file: {file_path}")
        text = _read_text(file_path)
        if text.strip():
            documents.append(LoadedDocument(path=file_path, content=text))
        else:
            logger.debug(f"No content extracted from file: {file_path}. Skipping.")
    
    # print(f"Loaded {len(documents)} documents.")
    logger.info(f"Loaded {len(documents)} documents.")
    return documents
        
# Reads the text content from a file at the given path.
def _read_text(path: Path)->str:
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md"}:
            return path.read_text(encoding="utf-8")
        elif suffix == ".pdf":
            reader = PdfReader(str(path))
            pages = [page.extract_text() or "" for page in reader.pages]
            return "\n".join(pages)
        elif suffix == ".pptx":
            presentation = Presentation(str(path))
            texts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts)
        else:
            logger.warning(f"Unsupported file type: {path}")
            return ""
    except Exception as e:
        logger.error(f"Error reading file {path}: {e}")
    return ""

# Splits the given text into chunks of specified size with overlap.
def chunk_text(text: str, *, chunk_size: int, chunk_overlap: int) -> list[str]:
    words = text.split()
    if not words:
        return []
    if chunk_size <= 0:
        raise ValueError("chunk_size must be greater than 0")
    if chunk_overlap >= chunk_size:
        raise ValueError("chunk_overlap must be smaller than chunk_size")

    step = chunk_size - chunk_overlap
    chunks: list[str] = []
    for start in range(0, len(words), step):
        end = min(start + chunk_size, len(words))
        chunk_words = words[start:end]
        chunks.append(" ".join(chunk_words))
        if end == len(words):
            break

    # print(chunks)
    return chunks

# Builds document chunks from loaded documents.
def build_chunks(documents: Sequence[LoadedDocument], settings: Settings) -> list[DocumentChunk]:
    active_settings = settings 
    chunks: list[DocumentChunk] = []
    for doc in documents:
        relative = doc.path.relative_to(active_settings.docs_path)
        document_id = relative.as_posix()
        safe_document_id = document_id.replace("/", "__")
        text_chunks = chunk_text(
            doc.content,
            chunk_size=active_settings.chunk_size,
            chunk_overlap=active_settings.chunk_overlap,
        )

        for index, chunk_content in enumerate(text_chunks):
            chunk_id = f"{safe_document_id}__chunk_{index}"
            chunks.append(
                DocumentChunk(
                    chunk_id=chunk_id,
                    document_id=document_id,
                    source_path=str(relative),
                    chunk_index=index,
                    content=chunk_content,
                )
            )
    # print(chunks)
    logger.info("Built %d chunks from %d documents", len(chunks), len(documents))
    # print(f"Built {len(chunks)} chunks from {len(documents)} documents")
    return chunks 
    
# Yields successive n-sized chunks from iterable.
def _batched(iterable: Sequence[DocumentChunk], batch_size: int) -> Iterator[Sequence[DocumentChunk]]:
        for index in range(0, len(iterable), batch_size):
            yield iterable[index : index + batch_size]

# Main ingestion function
def ingest_documents(settings: Settings | None = None) -> int:

    active_settings = settings or get_settings()

    # Load Documents
    documents = load_documents(active_settings.docs_path)
    if not documents:
        logger.info("No documents found")
        return 0

    # print("Chunking documents...")
    logger.info("Chunking documents...")

    # Chunking the loaded documents
    chunks = build_chunks(documents, active_settings)
    if not chunks:
        # print("No chunks created from documents.")
        logger.warning("Documents yielded no content after chunkting")
        return 0

    # Load API Key
    openai_api_key = active_settings.openai_api_key
    if not openai_api_key:
        logger.error("OPENAI_API_KEY environment variable not set.")
        raise RuntimeError("OPENAI_API_KEY environment variable not set.")

    client = OpenAIClient(
        OpenAIClientConfig(
            api_key=openai_api_key,
            embed_model=active_settings.embed_model,
            timeout=active_settings.openai_timeout,
            max_retries=active_settings.max_embed_retries,
        )
    )

    # Load the Vector Store
    store = get_vector_store()

    total_chunks = len(chunks)
    batch_size = int(os.getenv("EMBED_BATCH_SIZE", DEFAULT_EMBED_BATCH_SIZE))
    batch_size = max(1, batch_size)

    logger.info("Embedding %s chunks in batches of %s", total_chunks, batch_size)

    processed = 0
    # max_retries = int(os.getenv("EMBED_MAX_RETRIES", 3))
    # backoff_base = float(os.getenv("EMBED_BACKOFF_BASE", 1.0))  # seconds

    # Process chunks in batches
    for batch in _batched(chunks, batch_size):
        texts = [chunk.content for chunk in batch]

        # print(texts)

        embeddings = client.embed_texts([chunk.content for chunk in batch])
        
        # print(embeddings)

        store.upsert(batch, embeddings)

        logger.debug("Embedded %s items in current batch", len(embeddings))
        # print(f"Embedded {len(embeddings)} items in current batch")

        processed += len(batch)
        percent = math.floor((processed / total_chunks) * 100)
        logger.info("Indexed %s/%s chunks (%s%%)", processed, total_chunks, percent)

    logger.info("Completed ingestion for %s documents", len(documents))
    # print(f"Completed ingestion for {len(documents)} documents")
    return total_chunks

def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Ingest project documents into Chroma")
    parser.add_argument(
        "--docs",
        type=Path,
        help="Override docs directory (defaults to DOCS_PATH or docs/)",
    )
    parser.add_argument(
        "--batch-size",
        type=int,
        help="Override embedding batch size",
    )
    return parser.parse_args()

def main() -> None:
    args = _parse_args()
    settings = get_settings()
    overrides: dict[str, object] = {}
    if args.docs:
        overrides["docs_path"] = args.docs
    if args.batch_size:
        os.environ["EMBED_BATCH_SIZE"] = str(args.batch_size)
    if overrides:
        settings = replace(settings, **overrides)
    ingest_documents(settings)


if __name__ == "__main__":
    main()